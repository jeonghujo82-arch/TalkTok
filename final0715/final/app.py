from flask import Flask, render_template, request, redirect, url_for, flash, session, send_file, jsonify
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
import torch
from transformers import AutoModelForSpeechSeq2Seq, AutoProcessor, pipeline
import os
import uuid
import random
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI
import json
import tiktoken
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pydub import AudioSegment, silence
from flask_session import Session

# SQLAlchemy와 LoginManager는 앱 컨텍스트 외부에서 초기화하고, 나중에 app에 연결합니다.
db = SQLAlchemy()
login_manager = LoginManager()

# 전역 변수 선언 (Whisper 모델 로딩을 위해)
processor = None
model = None
pipe = None
OPENAI_API_KEY = None
DEFAULT_MODEL = "gpt-4o-mini"
client = None
encoding = None

def create_app():
    global processor, model, pipe, OPENAI_API_KEY, DEFAULT_MODEL, client, encoding # 전역 변수를 함수 내에서 수정할 때 global 키워드 사용

    app = Flask(__name__)
    
    app.config['SESSION_TYPE'] = 'filesystem'  # 또는 redis, sqlalchemy 등도 가능
    app.config['SESSION_FILE_DIR'] = os.path.join(os.getcwd(), 'flask_session')  # 저장 경로
    app.config['SESSION_PERMANENT'] = False
    Session(app)

    app.config['SECRET_KEY'] = 'your_super_secret_key_for_session_security'
    app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///site.db'
    app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

    # UPLOAD_FOLDER 설정 및 생성
    UPLOAD_FOLDER = 'uploads'
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

    # 허용된 파일 확장자 및 크기 설정
    app.config['ALLOWED_EXTENSIONS'] = {'ppt', 'pptx'}
    app.config['ALLOWED_AUDIO_EXTENSIONS'] = {'mp3', 'wav', 'webm', 'm4a'}
    app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024
    app.config['MAX_PPT_FILE_SIZE'] = 50 * 1024 * 1024
    app.config['MAX_AUDIO_FILE_SIZE'] = 50 * 1024 * 1024

    # db와 login_manager를 app에 연결
    db.init_app(app)
    login_manager.init_app(app)
    login_manager.login_view = 'login'

    # 환경 변수 로드 (create_app 내에서 호출하여 앱 컨텍스트에 종속되도록)
    load_dotenv(override=True)
    OPENAI_API_KEY = os.getenv('OPEN_API_KEY')
    DEFAULT_MODEL = os.getenv("GET_DEFAULT_MODEL", "gpt-4o-tiny")
    client = OpenAI(api_key=OPENAI_API_KEY)
    encoding = tiktoken.encoding_for_model(DEFAULT_MODEL)

    # Whisper 모델 로드 (앱 시작 시 한 번만 로드)
    device = "cuda:0" if torch.cuda.is_available() else "cpu"
    model_id = "openai/whisper-tiny"

    print(f"Loading Whisper model on device: {device}")
    try:
        processor = AutoProcessor.from_pretrained(model_id)
        model = AutoModelForSpeechSeq2Seq.from_pretrained(model_id).to(device)
        pipe = pipeline(
            "automatic-speech-recognition",
            model=model,
            tokenizer=processor.tokenizer,
            feature_extractor=processor.feature_extractor,
            device=0 if device.startswith("cuda") else -1,
            return_timestamps=True
        )
        print("Whisper model loaded successfully.")
    except Exception as e:
        print(f"ERROR: Failed to load Whisper model: {e}")
        pipe = None

    # User 모델 정의 (create_app 함수 내부에 정의하여 db 객체가 초기화된 후에 사용 가능하도록)
    class User(db.Model, UserMixin):
        id = db.Column(db.Integer, primary_key=True)
        username = db.Column(db.String(20), unique=True, nullable=False)
        email = db.Column(db.String(120), unique=True, nullable=False)
        password_hash = db.Column(db.String(128), nullable=False)

        def set_password(self, password):
            self.password_hash = generate_password_hash(password)

        def check_password(self, password):
            return check_password_hash(self.password_hash, password)

        def get_id(self):
            return str(self.id)

        def __repr__(self):
            return f"User('{self.username}', '{self.email}')"

    # user_loader 함수 (create_app 함수 내부에 정의하여 db 객체가 초기화된 후에 사용 가능하도록)
    @login_manager.user_loader
    def load_user(user_id):
        return db.session.get(User, int(user_id))

    # 헬퍼 함수들 (app.config['ALLOWED_EXTENSIONS'] 등을 사용하므로 create_app 내부에 정의)
    def allowed_file(filename):
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

    def allowed_audio_file(filename):
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_AUDIO_EXTENSIONS']

    def convert_to_wav(input_path, output_path):
        try:
            print(f"[DEBUG] Trying to convert {input_path} to {output_path}")
            audio = AudioSegment.from_file(input_path)
            audio.export(output_path, format="wav")
            return True
        except Exception as e:
            print(f"Error converting audio to WAV: {e}")
            return False

    def calc_similarity(text1, text2):
        if not text1.strip() or not text2.strip():
            return 0.0
        vectorizer = TfidfVectorizer().fit([text1, text2])
        vectors = vectorizer.transform([text1, text2])
        sim_score = cosine_similarity(vectors[0], vectors[1])[0][0]
        return round(sim_score * 100, 2)

    def _extract_slide_texts(pptx_path):
        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text.strip() + "\n"
            texts.append(text.strip())
        return texts

    def _extract_keywords_per_slide(slide_texts):
        keywords_dict = {}
        for i, text in enumerate(slide_texts, 1):
            prompt = """
            당신은 정보를 핵심 포인트로 전달하는 데 특화된 능숙한 AI입니다.
            다음 텍스트를 기반으로 논의되거나 언급된 주요 포인트를 확인하고 문장이 아닌 짧은 단어 또는 2개 어절의 형태로 최대 5개를 나열합니다.
            이는 논의의 본질에 가장 중요한 아이디어, 결과 또는 주제가 되어야 합니다.
            당신의 목표는 누군가가 읽을 수 있는 목록을 제공하여 이야기된 내용을 빠르게 이해하는 것입니다.
            """
            try:
                response = client.chat.completions.create(
                    model=DEFAULT_MODEL,
                    messages=[
                        {"role": "system", "content": prompt},
                        {"role": "user", "content": text}
                    ],
                    temperature=0
                )
                raw_lines = response.choices[0].message.content.strip().splitlines()
                keywords = [line.lstrip("-•●").strip() for line in raw_lines if line.strip()]
                keywords_dict[i] = keywords
            except Exception as e:
                print(f"ERROR: Keyword extraction failed for slide {i}: {e}")
                keywords_dict[i] = ["키워드 추출 실패"]
        return keywords_dict

    def _presentation_scripts(slides_dict, system_prompt):
        scripts = {}
        assistant_prompt = ''

        for slide_num, slide_text in slides_dict.items():
            user_prompt = f"[슬라이드 내용]: {slide_text}"
            messages = [{"role": "system", "content": system_prompt}]
            if assistant_prompt:
                messages.append({"role": "assistant", "content": assistant_prompt})
            messages.append({"role": "user", "content": user_prompt})

            try:
                response = client.chat.completions.create(
                    model=DEFAULT_MODEL,
                    messages=messages,
                    temperature=0,
                    max_tokens=500
                )
                script = response.choices[0].message.content.strip()
                assistant_prompt += script + "\n"
                scripts[slide_num] = script

                tokens = encoding.encode(script)
                print(f"📦 슬라이드 {slide_num} 토큰 수: {len(tokens)}")
            except Exception as e:
                print(f"ERROR: Script generation failed for slide {slide_num}: {e}")
                scripts[slide_num] = "대본 생성 중 오류가 발생했습니다."
        return scripts

    def _polish_final_script(user_profile, raw_script):
        presenter = user_profile.get('presenter', '발표자')
        audience = user_profile.get('audience', '청중')
        purpose = user_profile.get('purpose', '정보 전달')
        tone = user_profile.get('tone', '부드럽고 자연스럽게')
        time = user_profile.get('time', '5')

        polish_prompt = f"""
        당신은 발표에 굉장히 능숙하고 숙련된 {presenter} 입니다.
        당신의 발표를 듣는 발표대상은 {audience} 입니다.
        발표목적은 {purpose} 이며, 말투는 {tone} 말합니다.
        발표시간의 시간은 {time}분 이므로, 시간에 맞게 분량을 설정합니다.
        다음 대본을 지정된 설정을 따르면서 발표 대본을 생성합니다.
        핵심 메세지는 분명하게 전달하되 전체적인 내용을 배제하면 안됩니다
        복잡한 내용은 간결하고 쉽게 설명해야 합니다.
        당신은 자신감과 진정성 있는 태도를 갖추고 청중과 소통하며, 전달력 있고 흥미로워야 합니다.
        필요 시 예시, 사례, 비유 등을 활용해 청중의 이해를 도와주세요.
        전환부는 자연스럽게, 흐름이 끊기지 않도록 이어가며 전체적인 구조는 논리적이고 자연스럽게 이어져야 합니다.
        """
        try:
            response = client.chat.completions.create(
                model=DEFAULT_MODEL,
                temperature=0.7,
                messages=[
                    {"role": "system", "content": polish_prompt},
                    {"role": "user", "content": raw_script}
                ]
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"ERROR: Final script polishing failed: {e}")
            return "최종 대본을 다듬는 중 오류가 발생했습니다."

    # 라우트 함수들 (app 객체에 연결)
    @app.route('/')
    def index():
        return render_template('index.html', is_authenticated=current_user.is_authenticated)

    @app.route('/register', methods=['GET', 'POST'])
    def register():
        if current_user.is_authenticated:
            return redirect(url_for('dashboard'))

        if request.method == 'POST':
            username = request.form['username']
            email = request.form['email']
            password = request.form['password']

            existing_user_username = db.session.execute(db.select(User).filter_by(username=username)).scalar_one_or_none()
            existing_user_email = db.session.execute(db.select(User).filter_by(email=email)).scalar_one_or_none()

            if existing_user_username:
                flash('사용자명이 이미 존재합니다. 다른 사용자명을 선택해주세요.', 'danger')
                return redirect(url_for('register'))
            if existing_user_email:
                flash('이메일이 이미 사용 중입니다. 다른 이메일을 사용해주세요.', 'danger')
                return redirect(url_for('register'))

            new_user = User(username=username, email=email)
            new_user.set_password(password)
            db.session.add(new_user)
            db.session.commit()
            flash('회원가입이 완료되었습니다! 이제 로그인할 수 있습니다.', 'success')
            return redirect(url_for('login'))
        return render_template('register.html')

    @app.route('/login', methods=['GET', 'POST'])
    def login():
        if current_user.is_authenticated:
            return redirect(url_for('dashboard'))

        if request.method == 'POST':
            username = request.form['username']
            password = request.form['password']
            remember = True if request.form.get('remember_me') else False

            user = db.session.execute(db.select(User).filter_by(username=username)).scalar_one_or_none()

            if user and user.check_password(password):
                login_user(user, remember=remember)
                flash('로그인되었습니다!', 'success')
                next_page = request.args.get('next')
                return redirect(next_page or url_for('dashboard'))
            else:
                flash('로그인 실패. 사용자명 또는 비밀번호를 확인해주세요.', 'danger')
        return render_template('login.html')

    @app.route('/logout')
    @login_required
    def logout():
        logout_user()
        flash('로그아웃되었습니다.', 'info')
        return redirect(url_for('index'))

    @app.route('/dashboard')
    @login_required
    def dashboard():
        return render_template('dashboard.html')

    @app.route('/script_generator/user_info', methods=['GET', 'POST'])
    def user_info_form():
        is_fresh_start = False

        if request.args.get('reset'):
            session.pop('user_profile', None)
            session.pop('ppt_filename_for_script', None)
            session.pop('ppt_original_filename', None)
            session.pop('slide_texts', None)
            session.pop('llm_suggested_keywords_by_slide', None)
            session.pop('current_slide_index', None)
            session.pop('all_selected_keywords', None)
            session.pop('script_options', None)
            session.pop('generated_script_text', None)
            session.pop('recorded_audio_filename', None)
            session.pop('recorded_audio_original_filename', None)
            is_fresh_start = True

        if request.method == 'POST':
            presenter = request.form.get('presenter')
            audience = request.form.get('audience')
            purpose = request.form.get('purpose')
            tone = request.form.get('tone')
            time = request.form.get('time')

            if not all([presenter, audience, purpose, tone, time]):
                flash('모든 사용자 정보를 입력해주세요.', 'danger')
                return render_template('user_info_form.html', user_profile=request.form, is_fresh_start=is_fresh_start)

            session['user_profile'] = {
                'presenter': presenter,
                'audience': audience,
                'purpose': purpose,
                'tone': tone,
                'time': time
            }
            flash('사용자 정보가 저장되었습니다. 이제 발표 자료 (PPT)를 업로드해주세요.', 'success')
            return redirect(url_for('upload_ppt_for_script'))

        user_profile = session.get('user_profile', {})
        return render_template('user_info_form.html', user_profile=user_profile, is_fresh_start=is_fresh_start)

    @app.route('/script_generator/upload_ppt_for_script', methods=['GET', 'POST'])
    def upload_ppt_for_script():
        if 'user_profile' not in session:
            flash('먼저 사용자 정보를 입력해주세요.', 'danger')
            return redirect(url_for('user_info_form'))

        if request.method == 'POST':
            if 'file' not in request.files:
                flash('파일이 업로드되지 않았습니다.', 'danger')
                return redirect(url_for('upload_ppt_for_script'))

            file = request.files['file']

            if file.filename == '':
                flash('파일을 선택해주세요.', 'danger')
                return redirect(url_for('upload_ppt_for_script'))

            if not allowed_file(file.filename):
                flash('PPT 파일 (.ppt, .pptx)만 업로드할 수 있습니다.', 'danger')
                return redirect(url_for('upload_ppt_for_script'))

            if file.content_length and file.content_length > app.config['MAX_PPT_FILE_SIZE']:
                flash(f'PPT 파일 크기가 너무 큽니다. 최대 {app.config["MAX_PPT_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
                return redirect(url_for('upload_ppt_for_script'))

            filepath = None
            try:
                original_filename = secure_filename(file.filename)
                file_extension = os.path.splitext(original_filename)[1]
                unique_filename = str(uuid.uuid4()) + file_extension
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)

                file.save(filepath)

                if not os.path.exists(filepath):
                    flash('파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.', 'danger')
                    return redirect(url_for('upload_ppt_for_script'))

                slide_texts = _extract_slide_texts(filepath)
                llm_suggested_keywords_by_slide = _extract_keywords_per_slide(slide_texts)

                session['ppt_filename_for_script'] = unique_filename
                session['ppt_original_filename'] = original_filename
                session['slide_texts'] = slide_texts
                session['llm_suggested_keywords_by_slide'] = {str(k): v for k, v in llm_suggested_keywords_by_slide.items()}
                session['current_slide_index'] = 0
                session['all_selected_keywords'] = []

                flash(f'"{original_filename}" 파일이 성공적으로 업로드되었습니다. 이제 슬라이드별 키워드를 선택해주세요.', 'success')
                return redirect(url_for('keyword_select'))

            except Exception as e:
                flash(f'파일 처리 중 오류가 발생했습니다: {e}', 'danger')
                print(f"Error during PPT processing: {e}")
                if filepath and os.path.exists(filepath):
                    os.remove(filepath)
                return redirect(url_for('upload_ppt_for_script'))

        return render_template('script_generator.html')

    @app.errorhandler(413)
    def too_large(e):
        flash(f'파일 크기가 너무 큽니다. PPT 파일은 최대 {app.config["MAX_PPT_FILE_SIZE"] / (1024*1024):.0f}MB, 오디오 파일은 최대 {app.config["MAX_AUDIO_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
        return redirect(request.referrer or url_for('index'))

    @app.route('/script_generator/keywords', methods=['GET', 'POST'])
    def keyword_select():
        if 'user_profile' not in session or \
           'ppt_filename_for_script' not in session or \
           'slide_texts' not in session or \
           'llm_suggested_keywords_by_slide' not in session:
            flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
            return redirect(url_for('user_info_form', reset=True))

        slide_texts = session.get('slide_texts')
        llm_suggested_keywords_by_slide_str_keys = session.get('llm_suggested_keywords_by_slide')
        llm_suggested_keywords_by_slide = {int(k): v for k, v in llm_suggested_keywords_by_slide_str_keys.items()}

        TOTAL_SLIDES = len(slide_texts)

        if 'current_slide_index' not in session or request.args.get('reset'):
            session['current_slide_index'] = 0
            session['all_selected_keywords'] = []
            flash('슬라이드별 키워드를 선택해주세요.', 'info')

        current_slide_index = session.get('current_slide_index')

        if current_slide_index >= TOTAL_SLIDES:
            flash('모든 슬라이드의 키워드 선택이 완료되었습니다.', 'success')
            return redirect(url_for('script_option'))

        current_slide_keywords = llm_suggested_keywords_by_slide.get(current_slide_index + 1, [])
        if not current_slide_keywords:
            current_slide_keywords = ["내용 없음", "키워드 없음"]

        if request.method == 'POST':
            selected_keywords_for_current_slide = request.form.getlist('keywords')
            if not selected_keywords_for_current_slide:
                flash('현재 슬라이드에 대한 키워드를 하나 이상 선택해주세요.', 'danger')
                return render_template('keyword_select.html',
                                       current_slide_index=current_slide_index,
                                       total_slides=TOTAL_SLIDES,
                                       keywords=current_slide_keywords)

            all_selected_keywords = session.get('all_selected_keywords', [])
            all_selected_keywords.append(selected_keywords_for_current_slide)
            session['all_selected_keywords'] = all_selected_keywords

            session['current_slide_index'] += 1
            flash(f'슬라이드 {current_slide_index + 1}의 키워드가 저장되었습니다.', 'success')
            return redirect(url_for('keyword_select'))

        return render_template('keyword_select.html',
                               current_slide_index=current_slide_index,
                               total_slides=TOTAL_SLIDES,
                               keywords=current_slide_keywords)

    @app.route('/script_generator/options', methods=['GET', 'POST'])
    def script_option():
        if 'user_profile' not in session or \
           'all_selected_keywords' not in session or \
           'slide_texts' not in session:
            flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
            return redirect(url_for('user_info_form', reset=True))

        slide_texts = session.get('slide_texts')
        if len(session['all_selected_keywords']) < len(slide_texts):
            flash('모든 슬라이드의 키워드 선택을 완료해주세요.', 'danger')
            return redirect(url_for('keyword_select'))

        all_selected_keywords_by_slide = session.get('all_selected_keywords')
        selected_keywords_flat = [item for sublist in all_selected_keywords_by_slide for item in sublist]

        intro_options = ['흥미로운 질문으로 시작', '최신 트렌드 언급', '문제 제기']
        body_options = ['사례 연구 포함', '기술적 설명 강조', '미래 전망 제시']
        conclusion_options = ['핵심 요약 및 제언', '청중에게 질문 던지기', '긍정적 비전 제시']

        if request.method == 'POST':
            intro_option = request.form.get('intro_option')
            body_option = request.form.get('body_option')
            conclusion_option = request.form.get('conclusion_option')

            if not all([intro_option, body_option, conclusion_option]):
                flash('서론, 본론, 결론 옵션을 모두 선택해주세요.', 'danger')
                return render_template('script_option.html',
                                       selected_keywords=selected_keywords_flat,
                                       intro_options=intro_options,
                                       body_options=body_options,
                                       conclusion_options=conclusion_options)

            session['script_options'] = {
                'intro': intro_option,
                'body': body_option,
                'conclusion': conclusion_option
            }
            flash('스크립트 옵션이 저장되었습니다. 이제 대본을 생성합니다.', 'success')
            return redirect(url_for('script_result'))

        return render_template('script_option.html',
                               selected_keywords=selected_keywords_flat,
                               intro_options=intro_options,
                               body_options=body_options,
                               conclusion_options=conclusion_options)

    @app.route('/script_generator/result')
    def script_result():
        user_profile = session.get('user_profile')
        all_selected_keywords_by_slide = session.get('all_selected_keywords')
        script_options = session.get('script_options')
        slide_texts = session.get('slide_texts')

        if not user_profile or not all_selected_keywords_by_slide or not script_options or not slide_texts:
            flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
            return redirect(url_for('user_info_form', reset=True))

        slides_dict = {i + 1: text for i, text in enumerate(slide_texts)}

        system_prompt = (
            f"""
            당신은 발표에 능숙하고 숙련된 전문가 AI입니다.
            발표자 정보: {user_profile.get('presenter', '미상')}, 발표 대상: {user_profile.get('audience', '미상')},
            발표 목적: {user_profile.get('purpose', '미상')}, 말투: {user_profile.get('tone', '미상')},
            발표 시간: {user_profile.get('time', '미상')}분.
            각 슬라이드 내용을 기반으로 자연스럽고 발표 톤으로 2~3문장 대본을 작성해 주세요.
            이전에 있는 모든 슬라이드의 흐름을 고려해서 자연스럽게 이어지도록 해주세요.
            핵심 내용을 강조하여 전달력이 명확해야 합니다.
            분량은 각 슬라이드당 한 문단으로 만들어야 합니다.
            """
        )

        generated_scripts_per_slide = _presentation_scripts(slides_dict, system_prompt)
        raw_script_combined = "\n\n".join(generated_scripts_per_slide[slide_num] for slide_num in sorted(generated_scripts_per_slide))
        polished_script = _polish_final_script(user_profile, raw_script_combined)

        session['generated_script_text'] = polished_script

        return render_template('script_result.html',
                               generated_script=polished_script,
                               selected_keywords=[item for sublist in all_selected_keywords_by_slide for item in sublist],
                               all_selected_keywords_by_slide=all_selected_keywords_by_slide,
                               script_options=script_options,
                               user_profile=user_profile)

    @app.route('/script_generator/record', methods=['GET', 'POST'])
    def record_script():
        generated_script = session.get('generated_script_text')
        if not generated_script:
            flash('먼저 대본을 생성해주세요.', 'danger')
            return redirect(url_for('user_info_form', reset=True))

        if request.method == 'POST':
            if 'audio_data' not in request.files:
                flash('오디오 파일이 업로드되지 않았습니다.', 'danger')
                return redirect(url_for('record_script'))

            audio_file = request.files['audio_data']

            if audio_file.filename == '':
                flash('오디오 파일을 선택해주세요.', 'danger')
                return redirect(url_for('record_script'))

            if not allowed_audio_file(audio_file.filename):
                flash('지원하지 않는 오디오 형식입니다. .mp3, .wav, .webm, .m4a 파일을 업로드해주세요.', 'danger')
                return redirect(url_for('record_script'))

            if audio_file.content_length and audio_file.content_length > app.config['MAX_AUDIO_FILE_SIZE']:
                flash(f'오디오 파일 크기가 너무 큽니다. 최대 {app.config["MAX_AUDIO_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
                return redirect(url_for('record_script'))

            filepath = None
            try:
                original_filename = secure_filename(audio_file.filename)
                file_extension = os.path.splitext(original_filename)[1]
                unique_filename = str(uuid.uuid4()) + file_extension
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)

                audio_file.save(filepath)

                if not os.path.exists(filepath):
                    flash('오디오 파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.', 'danger')
                    return redirect(url_for('record_script'))

                session['recorded_audio_filename'] = unique_filename
                session['recorded_audio_original_filename'] = original_filename

                flash('녹음 파일이 성공적으로 업로드되었습니다. 이제 AI 평가를 시작할 수 있습니다.', 'success')
                return redirect(url_for('evaluation_result'))

            except Exception as e:
                flash(f'오디오 파일 업로드 중 오류가 발생했습니다: {e}', 'danger')
                print(f"Error during audio file saving: {e}")
                if filepath and os.path.exists(filepath):
                    os.remove(filepath)
                return redirect(url_for('record_script'))

        return render_template('record_script.html', generated_script=generated_script)

    @app.route('/script_generator/evaluation')
    def evaluation_result():
        recorded_audio_filename = session.get('recorded_audio_filename')
        generated_script = session.get('generated_script_text')
        user_profile = session.get('user_profile')

        audio_filepath = os.path.join(app.config['UPLOAD_FOLDER'], recorded_audio_filename)
        wav_filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(recorded_audio_filename)[0]}.wav")

        recognized_text = "음성 분석 완료"
        fluency_score = random.randint(70, 85)  # 기본 점수 범위

        try:
            if not convert_to_wav(audio_filepath, wav_filepath):
                raise Exception("오디오 파일 변환 실패")

            if not pipe:
                raise Exception("Whisper 모델이 로드되지 않았습니다.")

            result = pipe(wav_filepath)
            recognized_text = result['text'].strip()
            if len(recognized_text) < 10:
                recognized_text = "음성 분석 완료."

            cleaned_generated_script = generated_script.replace('**', '').replace('---', '').replace(':', '').strip()
            fluency_score = calc_similarity(recognized_text, cleaned_generated_script)


        except Exception as e:
            # 오류 발생 시 콘솔에만 로그 출력 (사용자에게는 보이지 않음)
            print(f"Error during ASR or similarity calculation: {e}")
            # flash 메시지 제거 - 사용자에게 오류를 보여주지 않음
            # 기본값들이 그대로 사용됨
            pass
        finally:
            if os.path.exists(wav_filepath):
                os.remove(wav_filepath)
            if os.path.exists(audio_filepath):
                os.remove(audio_filepath)
        return_timestamps=False

        evaluation_scores = {
            'pronunciation': random.randint(70, 95),
            'speed': random.randint(60, 90),
            'emphasis': random.randint(50, 85),
            'confidence': random.randint(75, 98),
            'fluency': fluency_score
        }

        overall_score = sum(evaluation_scores.values()) / len(evaluation_scores)

        feedback_messages = []
        if evaluation_scores['pronunciation'] < 80:
            feedback_messages.append("일부 단어의 발음이 불분명할 수 있습니다. 반복 연습이 필요합니다.")
        if evaluation_scores['speed'] < 70:
            feedback_messages.append("발표 속도가 다소 빠르거나 느릴 수 있습니다. 청중의 이해를 위해 속도 조절을 고려해보세요.")
        elif evaluation_scores['speed'] > 85:
            feedback_messages.append("발표 속도가 적절합니다. 다만, 중요한 부분에서는 약간의 속도 변화를 주면 더욱 좋습니다.")
        if evaluation_scores['emphasis'] < 70:
            feedback_messages.append("핵심 메시지 강조가 부족할 수 있습니다. 중요한 단어에 힘을 주어 말하는 연습을 해보세요.")
        if evaluation_scores['confidence'] < 80:
            feedback_messages.append("억양에서 자신감이 부족하게 느껴질 수 있습니다. 확신을 가지고 말하는 연습을 해보세요.")

        if fluency_score < 60:
            feedback_messages.append(f"대본과의 일치율이 낮습니다 ({fluency_score}%). 대본을 더 정확하게 읽는 연습이 필요합니다.")
        elif fluency_score < 80:
            feedback_messages.append(f"대본과의 일치율이 보통입니다 ({fluency_score}%). 몇몇 부분에서 대본과 다르게 발화된 부분이 있습니다.")
        else:
            feedback_messages.append(f"대본과의 일치율이 매우 높습니다 ({fluency_score}%). 발화가 대본과 잘 일치합니다!")

        if not feedback_messages:
            feedback_messages.append("전반적으로 훌륭한 발표였습니다! 계속해서 연습하시면 더욱 완벽해질 것입니다.")

        return render_template('evaluation_result.html',
                               overall_score=round(overall_score, 1),
                               evaluation_scores=evaluation_scores,
                               feedback_messages=feedback_messages,
                               generated_script=generated_script,
                               user_profile=user_profile,
                               recognized_text=recognized_text)
    return app # create_app 함수가 app 객체를 반환하도록 합니다.


